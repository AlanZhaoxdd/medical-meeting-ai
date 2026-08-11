from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable

from app.schemas.export import PptDeckSpec
from app.services.export_bundle import AnalysisBundle, source_index


SLIDE_TYPE_TITLES = {
    "cover": "会议汇报",
    "agenda": "会议议程",
    "summary": "会议核心摘要",
    "topics": "主要议题与参会者观点",
    "viewpoints": "参会者观点",
    "cutoff_questions": "切点问题分析",
    "charts": "数据图表",
    "consensus": "共识、分歧与待确认事项",
    "actions": "行动项与下一步建议",
    "sources": "引用来源",
    "end": "Thanks",
}

CUSTOMER_TEMPLATE_PATH = (
    Path(__file__).resolve().parent.parent
    / "assets"
    / "ppt"
    / "conference_ai_customer_template.pptx"
)
CUSTOMER_TEMPLATE_SOURCE = "客户提供的 ConferenceAI 用户手册模板"


def slide_default_title(slide_type: str) -> str:
    return SLIDE_TYPE_TITLES.get(slide_type, slide_type)


def _strip_speaker(bullet: str) -> str:
    return re.sub(r"^\s*[\[（(]?[^\]）)]{0,12}[\]）)]?\s*[:：]\s*", "", bullet).strip()


def _wrap_title(value: str, max_chars: int = 18) -> str:
    value = " ".join(str(value or "").split())
    if len(value) <= max_chars:
        return value
    return f"{value[:max_chars]}\n{value[max_chars: max_chars * 2]}".strip()


def _iter_text_shapes(slide: Any) -> Iterable[Any]:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def _replace_text(shape: Any, text: str, *, font_name: str, font_size: float, color: Any, bold: bool = True) -> None:
    from pptx.util import Pt

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    try:
        run._r.get_or_add_rPr().set("lang", "zh-CN")
        run._r.get_or_add_rPr().set("altLang", "en-US")
        run.font._rPr.set(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
            font_name,
        )
    except Exception:
        pass


def _remove_slide(presentation: Any, slide: Any) -> None:
    slide_id = next(
        (
            item
            for item in presentation.slides._sldIdLst
            if presentation.part.rels[item.rId].target_part is slide.part
        ),
        None,
    )
    if slide_id is None:
        return
    presentation.part.drop_rel(slide_id.rId)
    presentation.slides._sldIdLst.remove(slide_id)


def _find_layout(presentation: Any, name: str) -> Any:
    for layout in presentation.slide_layouts:
        if layout.name == name:
            return layout
    return presentation.slide_layouts[6]


def _chart_summary(chart_data: dict[str, Any], interpretation: str | None) -> list[str]:
    summary: list[str] = []
    if interpretation:
        summary.append(interpretation)
    denominator = chart_data.get("denominator") or {}
    valid = chart_data.get("valid_observation_count")
    if valid is None:
        valid = denominator.get("value") if isinstance(denominator, dict) else None
    if valid is not None:
        summary.append(
            f"样本人数：{valid}人"
            if chart_data.get("data_origin") == "demo"
            else f"有效样本：{valid}"
        )
    excluded = chart_data.get("excluded_observation_count")
    if excluded:
        summary.append(f"未纳入统计：{excluded} 条")
    return summary[:4]


def _build_deck(
    bundle: AnalysisBundle,
    spec: PptDeckSpec,
    *,
    include_charts: bool,
    include_references: bool,
    anonymous_attendees: bool,
    chart_images: dict[str, bytes],
    chart_interpretations: dict[str, str] | None = None,
    chart_data: dict[str, dict[str, Any]] | None = None,
    report_unit: str | None = None,
    presenter: str | None = None,
) -> Any:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    if not CUSTOMER_TEMPLATE_PATH.is_file():
        raise FileNotFoundError(f"PPT 模板文件不存在：{CUSTOMER_TEMPLATE_PATH}")

    presentation = Presentation(str(CUSTOMER_TEMPLATE_PATH))
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    navy = RGBColor(0x0B, 0x2A, 0x6C)
    teal = RGBColor(0x18, 0xA6, 0xD1)
    dark_teal = RGBColor(0x0E, 0x76, 0x99)
    text_color = RGBColor(0x22, 0x3E, 0x5C)
    muted = RGBColor(0x7E, 0x82, 0x86)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    interpretations = chart_interpretations or {}
    chart_specs = chart_data or {}

    def set_font(run: Any, size: float, color: Any, bold: bool = False, font_name: str = "Microsoft YaHei Light") -> None:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
        try:
            run._r.get_or_add_rPr().set("lang", "zh-CN")
            run._r.get_or_add_rPr().set("altLang", "en-US")
            run.font._rPr.set(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
                font_name,
            )
        except Exception:
            pass

    def add_textbox(
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        *,
        size: float,
        color: Any = text_color,
        bold: bool = False,
        font_name: str = "Microsoft YaHei Light",
        align: Any = PP_ALIGN.LEFT,
        valign: Any = MSO_ANCHOR.TOP,
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        set_font(run, size, color, bold=bold, font_name=font_name)
        return box

    def add_page_chrome(slide: Any, title: str, subtitle: str, page_number: int, total: int) -> None:
        add_textbox(slide, 0.48, 0.20, 11.9, 0.48, title, size=26, color=navy, bold=True)
        if subtitle:
            add_textbox(slide, 0.58, 0.73, 11.7, 0.34, subtitle, size=16, color=muted, bold=True)
        accent = slide.shapes.add_shape(1, Inches(0.48), Inches(1.20), Inches(0.82), Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = teal
        accent.line.fill.background()
        add_textbox(slide, 0.48, 7.10, 10.8, 0.22, deck_title, size=9.5, color=muted)
        add_textbox(slide, 12.15, 7.10, 0.68, 0.22, f"{page_number:02d}", size=10, color=dark_teal, bold=True, align=PP_ALIGN.RIGHT)

    def add_bullets(slide: Any, bullets: list[Any], *, left: float, top: float, width: float, height: float) -> None:
        usable = bullets[:8]
        if not usable:
            return
        row_height = min(0.72, max(0.44, height / max(len(usable), 1)))
        for index, bullet in enumerate(usable):
            if isinstance(bullet, str):
                text = bullet.strip()
            else:
                data = bullet.model_dump() if not isinstance(bullet, dict) else bullet
                text = str(data.get("text") or "").strip()
            if anonymous_attendees:
                text = _strip_speaker(text)
            if not text:
                continue
            add_textbox(
                slide,
                left,
                top + index * row_height,
                0.34,
                0.30,
                f"{index + 1:02d}",
                size=11,
                color=teal,
                bold=True,
            )
            line = slide.shapes.add_shape(1, Inches(left + 0.38), Inches(top + index * row_height + 0.02), Inches(0.035), Inches(row_height - 0.12))
            line.fill.solid()
            line.fill.fore_color.rgb = teal
            line.line.fill.background()
            add_textbox(slide, left + 0.56, top + index * row_height, width - 0.56, row_height, text, size=16, color=text_color)

    def add_sources(slide: Any) -> None:
        sources = bundle.sources[:24]
        if not sources:
            add_textbox(slide, 0.72, 1.65, 11.8, 0.6, "本次会议没有可展示的引用来源。", size=16, color=muted)
            return
        columns = [sources[:12], sources[12:24]]
        for col_index, items in enumerate(columns):
            left = 0.72 + col_index * 6.25
            for row_index, source in enumerate(items):
                title = str(source.get("title") or "来源")
                source_type = str(source.get("type") or "")
                value = f"[{source_index(source)}] {title}"
                if source_type:
                    value += f" · {source_type}"
                add_textbox(slide, left, 1.55 + row_index * 0.42, 5.75, 0.34, value, size=13, color=text_color)

    def populate_cover(slide: Any) -> None:
        text_shapes = sorted(_iter_text_shapes(slide), key=lambda item: item.top)
        if len(text_shapes) >= 2:
            _replace_text(text_shapes[0], _wrap_title(deck_title), font_name="Ubuntu", font_size=38, color=teal, bold=True)
            _replace_text(text_shapes[1], "ConferenceAI", font_name="Ubuntu", font_size=60, color=teal, bold=True)
        cover_meta = " · ".join(value for value in (deck_subtitle, report_unit, presenter) if value)
        if cover_meta:
            add_textbox(slide, 6.95, 5.78, 5.35, 0.46, cover_meta, size=13, color=white)
        add_textbox(slide, 6.95, 6.38, 5.35, 0.28, "AI 会议纪要分析报告", size=11, color=RGBColor(0x9F, 0xD6, 0xE6))

    def clear_placeholders(slide: Any) -> None:
        """Remove inherited empty placeholders so generated pages are self-contained."""
        for shape in list(slide.placeholders):
            element = shape._element
            parent = element.getparent()
            if parent is not None:
                parent.remove(element)

    def populate_agenda(slide: Any) -> None:
        text_shapes = sorted(_iter_text_shapes(slide), key=lambda item: (item.top, item.left))
        if not text_shapes:
            return
        agenda_title = next((shape for shape in text_shapes if shape.text.strip() == "Agenda"), text_shapes[0])
        _replace_text(agenda_title, "Agenda", font_name="Ubuntu", font_size=36, color=teal, bold=True)
        section_items = [
            ("01/", "会议概览"),
            ("02/", "核心内容"),
            ("03/", "切点图表"),
            ("04/", "结论与行动"),
        ]
        row_shapes = [shape for shape in text_shapes if shape is not agenda_title]
        for index, (number, label) in enumerate(section_items):
            if index * 2 + 1 >= len(row_shapes):
                break
            _replace_text(row_shapes[index * 2], number, font_name="Ubuntu", font_size=26, color=white, bold=True)
            _replace_text(row_shapes[index * 2 + 1], label, font_name="Microsoft YaHei Light", font_size=26, color=white, bold=True)

    def add_chart_content(slide: Any, slide_spec: Any, chart_id: str) -> None:
        data = chart_specs.get(chart_id, {})
        summary = _chart_summary(data, interpretations.get(chart_id))
        add_textbox(slide, 0.72, 1.55, 4.65, 0.35, "核心发现", size=18, color=dark_teal, bold=True)
        add_bullets(slide, summary or slide_spec.bullets, left=0.72, top=2.02, width=4.65, height=3.45)
        image = chart_images.get(chart_id)
        if image:
            slide.shapes.add_picture(BytesIO(image), Inches(5.55), Inches(1.48), width=Inches(7.15))
        unit = str(data.get("unit") or "")
        valid = data.get("valid_observation_count")
        excluded = data.get("excluded_observation_count") or 0
        metadata_parts: list[str] = []
        if unit:
            metadata_parts.append(f"单位：{unit}")
        if valid is not None:
            metadata_parts.append(
                f"样本人数：{valid}人"
                if data.get("data_origin") == "demo"
                else f"有效样本：{valid}"
            )
        if excluded:
            metadata_parts.append(f"排除：{excluded}")
        if metadata_parts:
            add_textbox(
                slide,
                5.55,
                6.10,
                7.15,
                0.30,
                " · ".join(metadata_parts),
                size=11,
                color=muted,
            )

    slides = list(spec.slides)
    total = len(slides)
    deck_title = spec.title or bundle.meeting.title
    deck_subtitle = spec.subtitle or ""
    if not deck_subtitle and bundle.meeting.starts_at:
        deck_subtitle = bundle.meeting.starts_at.strftime("%Y-%m-%d")
    if not deck_subtitle and bundle.meeting.organizer:
        deck_subtitle = bundle.meeting.organizer

    source_slides = list(presentation.slides)
    cover_source = source_slides[0] if source_slides else None
    agenda_source = source_slides[1] if len(source_slides) > 1 else None
    end_source = source_slides[14] if len(source_slides) > 14 else None
    cover_used = bool(slides and str(slides[0].type or "") == "cover" and cover_source is not None)
    agenda_used = bool(len(slides) > 1 and str(slides[1].type or "") == "agenda" and agenda_source is not None)
    end_used = bool(slides and str(slides[-1].type or "") == "end" and end_source is not None)
    retained = {
        next(
            item.rId
            for item in presentation.slides._sldIdLst
            if presentation.part.rels[item.rId].target_part is slide.part
        )
        for slide, used in (
            (cover_source, cover_used),
            (agenda_source, agenda_used),
            (end_source, end_used),
        )
        if slide is not None and used
    }
    for source_slide in source_slides:
        source_rid = next(
            (
                item.rId
                for item in presentation.slides._sldIdLst
                if presentation.part.rels[item.rId].target_part is source_slide.part
            ),
            None,
        )
        if source_rid not in retained:
            _remove_slide(presentation, source_slide)

    content_layout = _find_layout(presentation, "1_Title Only")
    built_slides: list[Any] = []
    for index, slide_spec in enumerate(slides, start=1):
        slide_type = str(slide_spec.type or "summary")
        title = slide_spec.title or slide_default_title(slide_type)
        if index == 1 and cover_used:
            slide = cover_source
            populate_cover(slide)
        elif index == 2 and agenda_used:
            slide = agenda_source
            populate_agenda(slide)
        elif index == total and end_used:
            slide = end_source
            end_shapes = sorted(_iter_text_shapes(slide), key=lambda item: item.top)
            if end_shapes:
                _replace_text(end_shapes[0], "Thanks", font_name="Ubuntu", font_size=60, color=teal, bold=True)
        else:
            slide = presentation.slides.add_slide(content_layout)
            clear_placeholders(slide)
            subtitle = " · ".join(value for value in (deck_subtitle, report_unit) if value) if index == 1 else ""
            add_page_chrome(slide, title, subtitle, index, total)
            if slide_type == "sources" and include_references:
                add_sources(slide)
            elif include_charts and slide_spec.chartIds:
                add_chart_content(slide, slide_spec, str(slide_spec.chartIds[0]))
            else:
                add_bullets(slide, slide_spec.bullets or [], left=0.72, top=1.58, width=11.7, height=4.95)
        notes = str(slide_spec.speakerNotes or "").strip()
        source_block = f"[Sources]\n- {CUSTOMER_TEMPLATE_SOURCE}"
        slide.notes_slide.notes_text_frame.text = f"{notes}\n\n{source_block}" if notes else source_block
        built_slides.append(slide)

    slide_ids = presentation.slides._sldIdLst
    ordered_ids = [
        next(
            item
            for item in slide_ids
            if presentation.part.rels[item.rId].target_part is slide.part
        )
        for slide in built_slides
    ]
    for slide_id in ordered_ids:
        slide_ids.remove(slide_id)
    for slide_id in ordered_ids:
        slide_ids.append(slide_id)

    return presentation


def render_ppt_bytes(
    bundle: AnalysisBundle,
    spec: PptDeckSpec,
    *,
    include_charts: bool,
    include_references: bool,
    anonymous_attendees: bool,
    chart_images: dict[str, bytes],
    chart_interpretations: dict[str, str] | None = None,
    chart_data: dict[str, dict[str, Any]] | None = None,
    report_unit: str | None = None,
    presenter: str | None = None,
) -> bytes:
    presentation = _build_deck(
        bundle,
        spec,
        include_charts=include_charts,
        include_references=include_references,
        anonymous_attendees=anonymous_attendees,
        chart_images=chart_images,
        chart_interpretations=chart_interpretations,
        chart_data=chart_data,
        report_unit=report_unit,
        presenter=presenter,
    )
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
