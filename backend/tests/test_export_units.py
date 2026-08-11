from __future__ import annotations

from datetime import datetime, timedelta, timezone
from io import BytesIO
from types import SimpleNamespace
from unittest.mock import AsyncMock
from uuid import uuid4

import pytest
from pptx import Presentation

from app.models.meeting import MeetingQuestion, MeetingQuestionType
from app.schemas.export import (
    PptBulletOut,
    PptDeckSpec,
    PptExportCreate,
    PptSlideOut,
    TextExportCreate,
)
from app.services.export_bundle import AnalysisBundle
from app.services.export_chart_service import (
    STANCE_LABELS,
    _aggregate_and_persist,
    _infer_speaker_from_text,
)
from app.services.export_model_clients import (
    ChartMentionItem,
    ChartMentionSet,
    ChartPlanResult,
    StanceItem,
)
from app.services.export_ppt import CUSTOMER_TEMPLATE_PATH, render_ppt_bytes
from app.services.export_text import (
    compose_export_sections,
    render_text_docx,
    render_text_pdf,
    split_minutes_sections,
)


def _meeting(**kwargs: object) -> SimpleNamespace:
    starts_at = datetime(2026, 8, 1, tzinfo=timezone.utc)
    values = {
        "id": uuid4(),
        "organization_id": uuid4(),
        "title": "医药顾问会",
        "starts_at": starts_at,
        "ends_at": starts_at + timedelta(hours=2),
        "location": "线上",
        "organizer": "市场部",
        "topic": "新品上市策略",
        "meeting_info": {"advisor_names": "张三, 李四", "internal_attendees": "王五"},
    }
    values.update(kwargs)
    return SimpleNamespace(**values)


def _question(question_type: MeetingQuestionType, content: str) -> MeetingQuestion:
    return MeetingQuestion(
        id=uuid4(),
        meeting_id=uuid4(),
        question_type=question_type,
        content=content,
        version=1,
        source="ai",
    )


def _bundle(
    *,
    modules: list[dict] | None = None,
    sources: list[dict] | None = None,
) -> AnalysisBundle:
    return AnalysisBundle(
        meeting=_meeting(),
        run=SimpleNamespace(
            source_version=3,
            modules=modules or [],
            sources=sources or [],
        ),
        questions=[],
        transcript_blocks=[],
    )


def test_split_minutes_sections_keeps_citation_indices() -> None:
    content = (
        "开篇\n"
        "## 会议总述\n总结 [1][2]\n"
        "## 行动项\n**跟进**（责任人：张三）[1]\n"
    )
    sections = split_minutes_sections(content)
    assert sections[0] == ("intro", "开篇", [])
    assert sections[1][0] == "会议总述"
    assert sections[1][2] == [1, 2]


def test_compose_export_sections_uses_only_confirmed_content() -> None:
    bundle = _bundle(
        modules=[
            {
                "id": "minutes",
                "title": "AI 通读纪要",
                "content": (
                    "**一、会议概述**\n会议达成一致 [1]\n"
                    "**二、分歧与焦虑**\n部分方案仍需确认 [1]\n"
                    "**三、循证数据解读**\n研究结果支持方案 [2]\n"
                    "**四、临床用药建议**\n建议根据人群调整 [2]\n"
                    "**五、专家共识**\n专家形成共识 [1]\n"
                    "**六、行动计划**\n第一，整理报告。[1]\n"
                ),
            }
        ],
        sources=[
            {"index": 1, "type": "transcript", "title": "转写片段", "snippet": "片段一"},
            {"index": 2, "type": "cutoff_question", "title": "切点问题", "snippet": "问题二"},
        ],
    )
    sections = compose_export_sections(
        bundle,
        show_attendee_names=True,
        include_references=True,
        include_citation_markers=True,
    )
    keys = [section.key for section in sections]
    assert keys[:6] == [
        "overview",
        "divergence",
        "evidence",
        "clinical",
        "consensus",
        "actions",
    ]
    assert "actions" in keys
    assert "sources" in keys
    assert all(
        section.citations
        for section in sections
        if section.key in {"summary", "cutoff", "actions"}
    )

    without_markers = compose_export_sections(
        bundle,
        show_attendee_names=True,
        include_references=False,
        include_citation_markers=False,
    )
    assert all("[" not in section.content for section in without_markers if section.content)


def test_text_export_request_supports_docx_and_pdf_without_optional_template_flags() -> None:
    for fmt in ("docx", "pdf"):
        payload = TextExportCreate(format=fmt)
        assert payload.format == fmt
        assert not hasattr(payload, "template")
        assert not hasattr(payload, "include_timestamps")


def test_render_text_docx_and_pdf_smoke() -> None:
    bundle = _bundle(
        modules=[
            {
                "id": "minutes",
                "title": "AI 通读纪要",
                "content": (
                    "**一、会议概述**\n核心结论 [1]\n"
                    "**六、行动计划**\n第一，跟进。[1]"
                ),
            }
        ],
        sources=[{"index": 1, "type": "transcript", "title": "转写", "snippet": "证据"}],
    )
    sections = compose_export_sections(
        bundle,
        show_attendee_names=True,
        include_references=True,
        include_citation_markers=True,
    )
    docx_bytes = render_text_docx(
        bundle,
        include_cover=True,
        sections=sections,
        include_references=True,
    )
    assert docx_bytes.startswith(b"PK")
    pdf_bytes = render_text_pdf(
        bundle,
        include_cover=True,
        sections=sections,
        include_references=True,
    )
    assert pdf_bytes.startswith(b"%PDF")


@pytest.mark.asyncio
async def test_bar_aggregation_deduplicates_speakers_and_validates_sources() -> None:
    meeting = _meeting()
    question = _question(MeetingQuestionType.CUT_POINT, "是否调整剂量？")
    bundle = AnalysisBundle(
        meeting=meeting,
        run=SimpleNamespace(
            source_version=3,
            modules=[],
            sources=[
                {
                    "index": 1,
                    "type": "transcript",
                    "title": "片段一",
                    "snippet": "我建议调整",
                    "speaker_name": "张三",
                },
                {
                    "index": 2,
                    "type": "transcript",
                    "title": "片段二",
                    "snippet": "我也同意",
                    "speaker_name": "张三",
                },
                {
                    "index": 3,
                    "type": "transcript",
                    "title": "片段三",
                    "snippet": "有条件支持",
                    "speaker_name": "李四",
                },
            ],
        ),
        questions=[question],
        transcript_blocks=[],
    )
    plan = ChartPlanResult(
        mentionSets=[
            ChartMentionSet(
                questionId=str(question.id),
                mentions=[
                    ChartMentionItem(speakerName="张三", sourceIds=["1", "2"]),
                    ChartMentionItem(speakerName="李四", sourceIds=["3"]),
                    # hallucinated speaker + source must be dropped
                    ChartMentionItem(speakerName="虚构者", sourceIds=["99"]),
                ],
            )
        ],
        stanceClassifications=[],
        planNote="",
    )
    session = AsyncMock()
    session.flush = AsyncMock()
    specs = await _aggregate_and_persist(
        session,
        bundle=bundle,
        plan=plan,
        chart_type="bar",
        target=question,
        metric="independent_speakers",
        source_by_id={
            "1": bundle.sources[0],
            "2": bundle.sources[1],
            "3": bundle.sources[2],
        },
        attendees=["张三", "李四", "王五"],
    )
    assert len(specs) == 1
    spec = specs[0]
    categories = spec.spec["categories"]
    assert len(categories) == 1
    assert categories[0]["value"] == 2
    assert len(categories[0]["evidence"]) == 3
    assert spec.valid is True


@pytest.mark.asyncio
async def test_pie_aggregation_is_mutually_exclusive_and_sums_to_denominator() -> None:
    meeting = _meeting()
    question = _question(MeetingQuestionType.CUT_POINT, "是否采用新方案？")
    bundle = AnalysisBundle(
        meeting=meeting,
        run=SimpleNamespace(
            source_version=3,
            modules=[],
            sources=[
                {
                    "index": 1,
                    "type": "transcript",
                    "title": "片段一",
                    "snippet": "支持",
                    "speaker_name": "张三",
                },
                {
                    "index": 2,
                    "type": "transcript",
                    "title": "片段二",
                    "snippet": "反对",
                    "speaker_name": "李四",
                },
            ],
        ),
        questions=[question],
        transcript_blocks=[],
    )
    plan = ChartPlanResult(
        mentionSets=[],
        stanceClassifications=[
            StanceItem(speakerName="张三", stance="SUPPORT", sourceIds=["1"]),
            StanceItem(speakerName="李四", stance="OPPOSE", sourceIds=["2"]),
        ],
        planNote="",
    )
    session = AsyncMock()
    session.flush = AsyncMock()
    specs = await _aggregate_and_persist(
        session,
        bundle=bundle,
        plan=plan,
        chart_type="pie",
        target=question,
        metric="stance_distribution",
        source_by_id={"1": bundle.sources[0], "2": bundle.sources[1]},
        attendees=["张三", "李四", "王五"],
    )
    spec = specs[0]
    total = sum(int(category["value"]) for category in spec.spec["categories"])
    assert total == 3
    assert spec.spec["denominator"]["value"] == 3
    labels = {category["label"] for category in spec.spec["categories"]}
    assert "未表态" in labels
    assert spec.valid is True


def test_infer_speaker_from_text_detects_inline_docx_speaker() -> None:
    attendees = ["洪天配", "崔瑾", "李冬梅", "吴红花"]
    assert _infer_speaker_from_text(
        "洪天配教授：尊敬的各位同道、各位朋友大家早上好。", attendees
    ) == "洪天配"
    # A later mention must not beat the speaker introducing the passage.
    assert _infer_speaker_from_text(
        "吴红花教授：在讨论之前，先回顾李冬梅教授提到的内容。", attendees
    ) == "吴红花"
    assert (
        _infer_speaker_from_text("开场致辞（09:00-09:05）： 洪天配教授：……", attendees)
        == "洪天配"
    )
    assert _infer_speaker_from_text("与任何参会者无关的段落。", attendees) is None


@pytest.mark.asyncio
async def test_bar_includes_open_questions_and_pie_accepts_open_target() -> None:
    meeting = _meeting()
    cut = _question(MeetingQuestionType.CUT_POINT, "是否采用新方案？")
    open_question = _question(MeetingQuestionType.OPEN_ENDED, "未来还需要补充哪些证据？")
    bundle = AnalysisBundle(
        meeting=meeting,
        run=SimpleNamespace(
            source_version=3,
            modules=[],
            sources=[
                {
                    "index": 1,
                    "type": "transcript",
                    "title": "片段一",
                    "snippet": "支持",
                    "speaker_name": "张三",
                },
                {
                    "index": 2,
                    "type": "transcript",
                    "title": "片段二",
                    "snippet": "反对",
                    "speaker_name": "李四",
                },
                {
                    "index": 3,
                    "type": "transcript",
                    "title": "片段三",
                    "snippet": "补充证据",
                    "speaker_name": "王五",
                },
            ],
        ),
        questions=[cut, open_question],
        transcript_blocks=[],
    )
    plan = ChartPlanResult(
        mentionSets=[
            ChartMentionSet(
                questionId=str(cut.id),
                mentions=[ChartMentionItem(speakerName="张三", sourceIds=["1"])],
            ),
            ChartMentionSet(
                questionId=str(open_question.id),
                mentions=[ChartMentionItem(speakerName="王五", sourceIds=["3"])],
            ),
        ],
        stanceClassifications=[
            StanceItem(speakerName="王五", stance="SUPPORT", sourceIds=["3"]),
        ],
        planNote="",
    )
    session = AsyncMock()
    session.flush = AsyncMock()
    source_by_id = {
        "1": bundle.sources[0],
        "2": bundle.sources[1],
        "3": bundle.sources[2],
    }
    bar_specs = await _aggregate_and_persist(
        session,
        bundle=bundle,
        plan=plan,
        chart_type="bar",
        target=cut,
        metric="independent_speakers",
        source_by_id=source_by_id,
        attendees=["张三", "李四", "王五"],
    )
    labels = [category["label"] for category in bar_specs[0].spec["categories"]]
    assert any(label.startswith("【切点】") for label in labels)
    assert any(label.startswith("【开放】") for label in labels)
    assert len(labels) == 2
    pie_specs = await _aggregate_and_persist(
        session,
        bundle=bundle,
        plan=plan,
        chart_type="pie",
        target=open_question,
        metric="stance_distribution",
        source_by_id=source_by_id,
        attendees=["张三", "李四", "王五"],
    )
    assert pie_specs[0].valid is True
    assert pie_specs[0].spec["title"].startswith("参会者立场分布：【开放】")


@pytest.mark.asyncio
async def test_bar_and_pie_specs_carry_deterministic_interpretation() -> None:
    meeting = _meeting()
    question = _question(MeetingQuestionType.CUT_POINT, "是否采用新方案？")
    bundle = AnalysisBundle(
        meeting=meeting,
        run=SimpleNamespace(
            source_version=3,
            modules=[],
            sources=[
                {
                    "index": 1,
                    "type": "transcript",
                    "title": "片段一",
                    "snippet": "支持",
                    "speaker_name": "张三",
                },
                {
                    "index": 2,
                    "type": "transcript",
                    "title": "片段二",
                    "snippet": "反对",
                    "speaker_name": "李四",
                },
            ],
        ),
        questions=[question],
        transcript_blocks=[],
    )
    plan = ChartPlanResult(
        mentionSets=[
            ChartMentionSet(
                questionId=str(question.id),
                mentions=[
                    ChartMentionItem(speakerName="张三", sourceIds=["1"]),
                    ChartMentionItem(speakerName="李四", sourceIds=["2"]),
                ],
            )
        ],
        stanceClassifications=[
            StanceItem(speakerName="张三", stance="SUPPORT", sourceIds=["1"]),
            StanceItem(speakerName="李四", stance="OPPOSE", sourceIds=["2"]),
        ],
        planNote="",
    )
    session = AsyncMock()
    session.flush = AsyncMock()
    bar_specs = await _aggregate_and_persist(
        session,
        bundle=bundle,
        plan=plan,
        chart_type="bar",
        target=question,
        metric="independent_speakers",
        source_by_id={"1": bundle.sources[0], "2": bundle.sources[1]},
        attendees=["张三", "李四", "王五"],
    )
    bar_interpretation = bar_specs[0].spec["interpretation"]
    assert "关注度最高" in bar_interpretation
    assert "独立参会者" in bar_interpretation
    assert "张三" in bar_interpretation or "李四" in bar_interpretation
    pie_specs = await _aggregate_and_persist(
        session,
        bundle=bundle,
        plan=plan,
        chart_type="pie",
        target=question,
        metric="stance_distribution",
        source_by_id={"1": bundle.sources[0], "2": bundle.sources[1]},
        attendees=["张三", "李四", "王五"],
    )
    pie_interpretation = pie_specs[0].spec["interpretation"]
    assert "明确支持" in pie_interpretation
    assert "未就该问题表态" in pie_interpretation


def test_ppt_deck_spec_validation_enforces_six_to_eight_slides() -> None:
    slides = [
        PptSlideOut(pageNumber=index, type="summary", title=f"页 {index}", bullets=[])
        for index in range(1, 6)
    ]
    with pytest.raises(ValueError):
        PptDeckSpec(title="汇报", slides=slides)


def test_ppt_export_request_accepts_chart_expanded_deck() -> None:
    slides = [
        PptSlideOut(pageNumber=index, type="summary", title=f"页 {index}", bullets=[])
        for index in range(1, 13)
    ]

    payload = PptExportCreate(slides=slides)

    assert payload.slides is not None
    assert len(payload.slides) == 12


def test_render_ppt_bytes_produces_editable_pptx() -> None:
    meeting = _meeting()
    bundle = SimpleNamespace(
        meeting=meeting,
        run=SimpleNamespace(source_version=3, modules=[]),
        questions=[],
        transcript_blocks=[],
        sources=[{"index": 1, "type": "transcript", "title": "转写", "snippet": "证据"}],
    )
    spec = PptDeckSpec(
        title="医药会议汇报",
        subtitle="2026-08-01",
        theme="formal",
        slides=[
            PptSlideOut(pageNumber=1, type="cover", title="封面", bullets=[]),
            PptSlideOut(
                pageNumber=2,
                type="summary",
                title="核心摘要",
                bullets=[PptBulletOut(text="张三：同意调整剂量", sourceIds=["1"])],
            ),
            PptSlideOut(
                pageNumber=3,
                type="consensus",
                title="共识",
                bullets=[PptBulletOut(text="达成一致", sourceIds=["1"])],
            ),
            PptSlideOut(
                pageNumber=4,
                type="actions",
                title="行动项",
                bullets=[PptBulletOut(text="整理报告", sourceIds=["1"])],
            ),
            PptSlideOut(pageNumber=5, type="sources", title="引用来源", bullets=[]),
            PptSlideOut(
                pageNumber=6,
                type="charts",
                title="数据图表",
                bullets=[PptBulletOut(text="覆盖度", sourceIds=[])],
                chartIds=["demo-chart"],
            ),
        ],
    )
    data = render_ppt_bytes(
        bundle,
        spec,
        include_charts=True,
        include_references=True,
        anonymous_attendees=True,
        chart_images={},
        chart_data={
            "demo-chart": {
                "data_origin": "demo",
                "unit": "%",
                "valid_observation_count": 116,
                "excluded_observation_count": 0,
                "denominator": {"name": "人数", "value": 116},
            }
        },
    )
    assert data.startswith(b"PK")
    assert len(data) > 10_000
    assert CUSTOMER_TEMPLATE_PATH.is_file()

    presentation = Presentation(BytesIO(data))
    assert len(presentation.slides) == 6
    cover_text = " ".join(
        shape.text
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "ConferenceAI" in cover_text
    assert "User Manual" not in cover_text
    assert "[Sources]" in presentation.slides[0].notes_slide.notes_text_frame.text
    chart_text = " ".join(
        shape.text
        for shape in presentation.slides[5].shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "样本人数：116人" in chart_text
    assert "统计口径" not in chart_text
    assert "独立参会者人数" not in chart_text
